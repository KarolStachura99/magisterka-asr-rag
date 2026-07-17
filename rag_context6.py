"""rag_context 6: hybrydowa ekstrakcja (unigramy KeyBERT + frazy noun_chunks,
formy powierzchniowe, filtr frekwencyjny). Tryby: plik PDF/PPTX lub --context-text."""
import os, sys, time, argparse
from collections import Counter
import spacy
from keybert import KeyBERT
from pypdf import PdfReader
from pptx import Presentation
try:
    from wordfreq import zipf_frequency
    HAS_WF = True
except ImportError:
    HAS_WF = False

print("Ładowanie spaCy...")
nlp = spacy.load("pl_core_news_md")
nlp.max_length = 2_000_000

def read_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        r = PdfReader(path)
        return " ".join(p.extract_text() for p in r.pages if p.extract_text())
    if ext == ".pptx":
        prs = Presentation(path)
        return " ".join(s.text for sl in prs.slides for s in sl.shapes if hasattr(s, "text"))
    raise ValueError(f"Nieobsługiwany format: {ext}")

def is_generic(word):
    """Filtr frekwencyjny: odrzuca słownictwo języka ogólnego (zipf>4.5)."""
    return HAS_WF and zipf_frequency(word.lower(), "pl") > 4.5

def extract_hybrid(raw_text, n_uni=20, n_phr=10):
    raw_text = " ".join(raw_text.split())   # normalizacja białych znaków z PDF
    doc = nlp(raw_text)
    # mapowanie lemat -> najczęstsza forma powierzchniowa
    surface = {}
    for t in doc:
        if t.is_alpha and len(t.text) > 2:
            surface.setdefault(t.lemma_.lower(), Counter())[t.text.lower()] += 1

# 1) FRAZY: wzorce POS (spaCy nie ma noun_chunks dla 'pl')
    chunks = Counter()
    tokens = list(doc)
    ok_pos = {"NOUN", "PROPN", "ADJ", "X"}
    for i in range(len(tokens)):
        for L in (2, 3):
            seq = tokens[i:i + L]
            if len(seq) < L:
                break
            if any(t.is_stop or t.is_punct or not t.is_alpha
                   or len(t.text) < 3 for t in seq):
                continue
            if not all(t.pos_ in ok_pos for t in seq):
                continue
            if not any(t.pos_ in {"NOUN", "PROPN", "X"} for t in seq):
                continue  # fraza musi zawierać rzeczownik
            if any(t.is_sent_start for t in seq[1:]):
                continue  # nie przecinamy granic zdań
            chunks[" ".join(t.text.lower() for t in seq)] += 1
    print("TOP FRAZY:", chunks.most_common(15))
    phrases = [p for p, c in chunks.most_common(40)
               if c >= 1 and not all(is_generic(w) for w in p.split())][:n_phr]

    # 2) UNIGRAMY: preprocessing jak rag_context 5 -> KeyBERT -> filtr -> formy powierzchniowe
    allowed = {"NOUN", "PROPN", "ADJ", "X"}
    cleaned = " ".join(
        (t.text if t.pos_ in {"PROPN", "X"} or t.text.isupper() else t.lemma_.lower())
        for t in doc
        if t.pos_ in allowed and not (t.is_stop or t.is_punct or len(t.text) < 3))
    kw = KeyBERT(model="paraphrase-multilingual-MiniLM-L12-v2")
    cands = kw.extract_keywords(cleaned, keyphrase_ngram_range=(1, 1),
                                top_n=40, use_mmr=True, diversity=0.5)
    unigrams = []
    for w, _ in cands:
        w = w.replace("_", " ")
        if is_generic(w) or any(w in p for p in phrases):
            continue
        best = surface.get(w.lower())
        unigrams.append(best.most_common(1)[0][0] if best else w)
        if len(unigrams) == n_uni:
            break
    return ", ".join(phrases + unigrams)

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--file", help="ścieżka do PDF/PPTX")
    ap.add_argument("--context-text", help="wolny opis sytuacji (bez ekstrakcji)")
    a = ap.parse_args()
    os.makedirs("extracted_dictionaries", exist_ok=True)
    t0 = time.time()
    if a.context_text:                       # tryb bez PDF: tekst idzie wprost
        wynik, name = a.context_text.strip(), "user_context"
    elif a.file:
        wynik = extract_hybrid(read_file(a.file))
        name = os.path.splitext(os.path.basename(a.file))[0]
    else:
        sys.exit("Podaj --file lub --context-text")
    out = f"extracted_dictionaries/extracted_dictionary_v6_{name}.txt"
    with open(out, "w", encoding="utf-8") as f:
        f.write(wynik)
    print(f"\nSłownik ({len(wynik.split(','))} pozycji):\n{wynik}")
    print(f"\n--- CZAS: {time.time()-t0:.2f} s ---\n[SUKCES] Zapisano: {out}")