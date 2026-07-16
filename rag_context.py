import os
import time
import spacy
from keybert import KeyBERT
from pypdf import PdfReader
from pptx import Presentation

print("Ładowanie modelu językowego spaCy...")
nlp = spacy.load("pl_core_news_md")

def extract_text_from_pdf(pdf_path: str) -> str:
    reader = PdfReader(pdf_path)
    return " ".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text_from_pptx(pptx_path: str) -> str:
    prs = Presentation(pptx_path)
    full_text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text += shape.text + " "
    return full_text

def preprocess_text_with_spacy(text: str) -> str:
    """
    Filtracja lingwistyczna (POS Tagging + Lemmatization).
    Przepuszcza tylko rzeczowniki (NOUN), nazwy własne (PROPN), przymiotniki (ADJ) 
    oras słowa obce/nieznane (X), które często są nazwami bibliotek IT.
    """
    nlp.max_length = 2000000 
    doc = nlp(text)
    
    allowed_pos = {'NOUN', 'PROPN', 'ADJ', 'X'}
    cleaned_tokens = []
    
    for token in doc:
        if token.is_stop or token.is_punct or len(token.text) < 3:
            continue
            
        if token.pos_ in allowed_pos:
            if token.pos_ in {'PROPN', 'X'} or token.text.isupper():
                cleaned_tokens.append(token.text)
            else:
                cleaned_tokens.append(token.lemma_.lower())
            
    return " ".join(cleaned_tokens)

def generate_keywords_keybert(file_path: str) -> str:
    _, ext = os.path.splitext(file_path)
    
    if ext.lower() == '.pdf':
        raw_text = extract_text_from_pdf(file_path)
    elif ext.lower() == '.pptx':
        raw_text = extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Nieobsługiwany format pliku: {ext}")
        
    print("Trwa lematyzacja i tagowanie POS (spaCy)...")
    doc_text = preprocess_text_with_spacy(raw_text)
    
    print("Trwa ekstrakcja wektorowa (KeyBERT)...")
    kw_model = KeyBERT(model='paraphrase-multilingual-MiniLM-L12-v2')
    
    keywords_with_scores = kw_model.extract_keywords(
        doc_text, 
        keyphrase_ngram_range=(1, 1), 
        stop_words=None,              
        top_n=30,                     
        use_mmr=True,                 
        diversity=0.5                 
    )
    
    keywords_list = [kw[0].replace("_", " ") for kw in keywords_with_scores]
    return ", ".join(keywords_list)

if __name__ == "__main__":
    sciezka_do_pliku = "C:/Users/Karol Stachura/Desktop/Projekt_Magisterski/pdf_files/Lab2 - Wstęp do klasyfikacji.pdf"
    
    nazwa_pliku = os.path.splitext(os.path.basename(sciezka_do_pliku))[0]
    
    output_dir = "extracted_dictionaries"
    os.makedirs(output_dir, exist_ok=True)
    
    plik_wyjsciowy = os.path.join(output_dir, f"extracted_dictionary_{nazwa_pliku}.txt")
    
    print(f"Uruchamiam Zaawansowany Potok RAG dla pliku: {nazwa_pliku}...")
    start_time = time.time()
    
    try:
        wynik = generate_keywords_keybert(sciezka_do_pliku)
        end_time = time.time()
        
        print(f"\nWygenerowany initial_prompt dla Whispera:\n{wynik}")
        print(f"\n--- CZAS OPERACJI EKSTRAKCJI: {end_time - start_time:.2f} sekund ---")
        
        with open(plik_wyjsciowy, "w", encoding="utf-8") as f:
            f.write(wynik)
        print(f"\n[SUKCES] Słownik został automatycznie zapisany do: {plik_wyjsciowy}")
        
    except Exception as e:
        print(f"\nBŁĄD PROCESU: {e}")